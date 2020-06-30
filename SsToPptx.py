import pyautogui
from pynput import keyboard
from pptx import Presentation
from pptx.util import Inches
import os,sys
from PIL import Image
from PIL.ImageQt import ImageQt
import time
from PyQt5 import QtWidgets,QtGui,QtCore
import threading

	
def resizeImage(image,size=(1920,1080),scale=1/3):
	im = image.resize((int(size[0]*scale),int(size[1]*scale)),Image.ANTIALIAS)
	return im

class Pencere(QtWidgets.QWidget):
	previews = []
	def __init_(self):
		QtWidgets.QWidget.__init__(self)
		QtCore.QThread.__init__(self)
	def run(self):
		self.home()
	def home(self):
		self.setWindowTitle("SstoPptx")
		self.h_box1 = QtWidgets.QHBoxLayout()
		self.v_box1 = QtWidgets.QVBoxLayout()
		self.v_box2 = QtWidgets.QVBoxLayout()

		#self.btn = QtWidgets.QPushButton("this")	
		#self.btn.clicked.connect(self.change)

		self.text = QtWidgets.QLabel("Latest screenshot")

		self.labelImage = QtWidgets.QLabel(self)	
		self.image = pyautogui.screenshot()
		self.image = resizeImage(self.image)
		self.qimage = ImageQt(self.image)
		self.pixmap = QtGui.QPixmap().fromImage(self.qimage)
		self.labelImage.setPixmap(self.pixmap)


		#self.v_box1.addWidget(self.btn)
		self.v_box2.addWidget(self.text)
		self.v_box2.addWidget(self.labelImage)		
		#self.h_box1.addLayout(self.v_box1)		
		self.h_box1.addLayout(self.v_box2)
		self.setLayout(self.h_box1)		
		
	def change(self,image,len):
		self.image = resizeImage(image)
		self.qimage = ImageQt(self.image)
		self.pixmap = QtGui.QPixmap().fromImage(self.qimage)		
		self.labelImage.setPixmap(self.pixmap)
		"""
		self.previews.append(QtWidgets.QLabel())
		self.previews[-1].setPixmap(self.pixmap)
		"""
		self.text.setText(str(len)+". screenshot (latest)")


class Slide:
	prs = Presentation()
	prs.slide_height = Inches(18/1.2)
	prs.slide_width = Inches(32/1.2)
	blank_slide_layout = prs.slide_layouts[6]
	title_slide_layout = prs.slide_layouts[0]
	Slides = []
	screenshots = []
	imageName = None
	left = Inches(0)
	top = Inches(0)
	width,height = 1920,1080
	width_,height_ = int(width/3) , int(height/3)	
	pptxName = ""
	tempPath = r"D:\TOBB ETU\SstoPptx\temp\\"

	def addText(self):
		title = input("Title\t: ")
		subtitle = input("Subtitle\t:")
		self.screenshots.append((title,subtitle))

	def createNewSlide(self,type):
		if(type == "title"):
			self.Slides.append(self.prs.slides.add_slide(self.title_slide_layout))
		elif(type == "image"):
			self.Slides.append(self.prs.slides.add_slide(self.blank_slide_layout))		


	def takeSS(self):	
		self.screenshots.append(pyautogui.screenshot())
		changePreview(self.screenshots[-1])


	def convertToPptx(self):
		t = time.localtime()
		timestamp = time.strftime('%b-%d-%Y_%H%M', t)
		self.pptxName = r"D:\TOBB ETU\SstoPptx" + r"\Slide  "+timestamp+".pptx"
		for i in range(len(self.screenshots)):			
			if(isinstance(self.screenshots[i],Image.Image)):				
				self.screenshots[i].save(self.tempPath + "ss"+str(i)+".png")##TODO: CLEAN THIS FOLDER BEFORE EXITING			
				self.createNewSlide("image")
				self.Slides[i].shapes.add_picture(self.tempPath + "ss"+str(i)+".png",self.left,self.top)
			elif(isinstance(self.screenshots[i],tuple)):
				self.createNewSlide("title")
				title = self.Slides[i].shapes.title
				subtitle = self.Slides[i].placeholders[1]
				title.text = self.screenshots[i][0]
				subtitle.text = self.screenshots[i][1]
				#self.Slides[i].shapes.title.text = self.screenshots[i][0]
				#self.Slides[i].placeholders[1].text = self.screenshots[i][1]	
		self.prs.save(self.pptxName)
		print("Saved to "+self.pptxName)

slide = Slide()
print("App Started ... \nListening for Key Inputs...\nScroll Lock\t-> TakeSS\nPause\t\t-> Convert\nEsc\t\t\t-> !!Exit!! ")
def on_press(key):
    try:    	
        #print('alphanumeric key {0} pressed'.format(key.char))
        pass
    except AttributeError:    	
        #print('special key {0} pressed'.format(key))
        pass

def on_release(key):
    #print('{0} released'.format(key))
    if key == keyboard.Key.esc:
    	if(input("Are you sure you want to quit without saving the file ? [Y/n]").lower() == 'y'):
    		return False
    	else:
    		print("Exiting aborted.\n")
    if key == keyboard.Key.pause:
        # Stop listener
        """  
        try:
        	slide.convertToPptx()
        	return False
        except:
        	print("!!! ERROR !!! File is safe.\t-->Check if powerPoint file is already opened and try again")
        """
        slide.convertToPptx()
        return False   
        
    if key == keyboard.Key.scroll_lock:
    	print(f"Scroll Lock key pressed : {len(slide.screenshots)+1}.Screenshot Taken")
    	#print('{0} released -> screenshot taken'.format(key))
    	slide.takeSS()

    if key == keyboard.Key.f9:
    	print("Title Slide created.\n")
    	slide.addText()

def changePreview(image):
	pencerem.change(image,len(slide.screenshots))



listener = keyboard.Listener(
	    on_press=on_press,
	    on_release=on_release)
listener.start()

app = QtWidgets.QApplication(sys.argv)
pencerem = Pencere()
pencerem.run()
#pencerem.home()	
pencerem.show()
sys.exit(app.exec_())

#main()

